#!/usr/bin/env python3
"""
🎬 视频 → 分镜头脚本 PPT 生成器
Video-to-Storyboard PPT Generator

输入一个视频链接，自动下载、场景分析、提取关键帧，
生成导演级分镜头脚本 PPT（黑白调 16:9 格式）。

依赖: yt-dlp, ffmpeg, python-pptx, Pillow
"""

import subprocess, os, sys, json, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

FONT = "Noto Sans SC"
OUTPUT_DIR = os.path.expanduser("~/Desktop")


def download_video(url):
    """下载视频并返回元数据"""
    print(f"📥 下载视频: {url}")
    r = subprocess.run(
        ["yt-dlp", "--dump-json", "--no-download", url],
        capture_output=True, text=True, timeout=30
    )
    if r.returncode != 0:
        raise RuntimeError(f"无法获取视频信息: {r.stderr[:200]}")
    meta = json.loads(r.stdout)

    r2 = subprocess.run(
        ["yt-dlp", "-f", "best", "-o", "temp_video.mp4", url],
        capture_output=True, text=True, timeout=120
    )
    if r2.returncode != 0:
        raise RuntimeError(f"下载失败: {r2.stderr[:200]}")

    print(f"  标题: {meta.get('title', '?')}")
    print(f"  时长: {meta.get('duration', '?')}秒")
    print(f"  分辨率: {meta.get('width', '?')}x{meta.get('height', '?')}")
    return meta


def detect_scenes():
    """场景检测，返回镜头切换时间点列表"""
    print("🔍 场景检测...")
    r = subprocess.run(
        ["ffmpeg", "-i", "temp_video.mp4",
         "-filter:v", "select='gt(scene,0.3)',showinfo",
         "-vsync", "vfr", "-q:v", "2", "-f", "null", "-"],
        capture_output=True, text=True, timeout=60
    )
    times = []
    for line in r.stderr.split('\n'):
        m = re.search(r'pts_time:([\d.]+)', line)
        if m:
            times.append(float(m.group(1)))

    # 加上0和总时长作为起止
    # 获取总时长
    r2 = subprocess.run(
        ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
         "-of", "csv=p=0", "temp_video.mp4"],
        capture_output=True, text=True, timeout=10
    )
    duration = float(r2.stdout.strip()) if r2.stdout.strip() else 30

    # 去重排序
    times = sorted(set([round(t, 2) for t in times]))
    print(f"  检测到 {len(times)} 个场景切换点")
    return times, duration


def build_shot_segments(scene_times, duration):
    """根据场景切换点构建镜头分段，返回 [(start, end, mid), ...]"""
    segments = []
    prev = 0.0
    for t in scene_times:
        if t - prev > 0.5:  # 过滤太短的片段
            mid = (prev + t) / 2
            segments.append((prev, t, mid))
            prev = t
    # 最后一段
    if duration - prev > 0.5:
        mid = (prev + duration) / 2
        segments.append((prev, duration, mid))
    return segments


def extract_frames(segments):
    """提取每镜中间帧"""
    os.makedirs("shots", exist_ok=True)
    print(f"🖼️ 提取 {len(segments)} 个关键帧...")
    for i, (start, end, mid) in enumerate(segments, 1):
        out = f"shots/shot_{i:02d}.jpg"
        r = subprocess.run(
            ["ffmpeg", "-y", "-ss", str(mid), "-i", "temp_video.mp4",
             "-vframes", "1", "-q:v", "2", out],
            capture_output=True, timeout=15
        )
        sz = os.path.getsize(out) // 1024 if os.path.exists(out) else 0
        print(f"  镜{i:02d}  @ {mid:.1f}s → {sz}KB")
    return len(segments)


def format_time(seconds):
    """秒 → 分:秒 格式"""
    m = int(seconds // 60)
    s = int(seconds % 60)
    return f"{m}:{s:02d}"


def generate_storyboard_data(meta, segments):
    """
    基于元数据和镜头分段构建分镜数据。
    用户需根据实际画面内容手动调整 scene/person/action 字段。
    """
    title = meta.get("title", "未命名视频")
    duration = meta.get("duration", 30)
    desc = meta.get("description", "") or ""
    tags = meta.get("tags", [])

    # 将描述文案按段落分割，分配到各个镜头
    paragraphs = [p.strip() for p in desc.split('\n') if p.strip()]
    # 过滤掉标签行
    script_lines = [p for p in paragraphs if not p.startswith('#')]

    shots = []
    for i, (start, end, mid) in enumerate(segments, 1):
        time_str = f"{format_time(start)}-{format_time(end)}"
        # 估算时长（秒）
        dur_sec = int(end - start)
        time_display = f"{dur_sec}s"

        # 从描述中取对应文案
        script = script_lines[i - 1] if i - 1 < len(script_lines) else "—"
        # 去掉引号
        script = script.strip('"').strip('"').strip()

        shots.append({
            "num": i,
            "type": "—",        # 用户需填写
            "time": time_display,
            "camera": "—",      # 用户需填写
            "scene": "",        # 用户需填写
            "person": "",       # 用户需填写
            "action": "",       # 用户需填写
            "script": script,
            "prop": "—",        # 用户需填写
            "title": title,
        })

    # 最后一段的文案如果没有，尝试用标题
    if shots and shots[-1]["script"] == "—" and title:
        shots[-1]["script"] = title

    return title, shots


def make_pptx(shots_data, title, output_path):
    """生成黑白调 16:9 PPT"""
    print(f"📊 生成 PPT...")
    W = Inches(13.333)
    H = Inches(7.5)

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    PAGE_BG = RGBColor(0xF7, 0xF7, 0xF7)
    BLACK = RGBColor(0x1A, 0x1A, 0x1A)
    DARK = RGBColor(0x33, 0x33, 0x33)
    MID = RGBColor(0x99, 0x99, 0x99)

    col_hdrs = ["镜号", "画面", "景别", "时长", "内容", "文案", "运镜", "道具"]
    col_ws = [0.55, 1.7, 0.85, 0.7, 3.7, 2.8, 0.8, 0.7]
    ncols = len(col_hdrs)
    TBL_L = Inches(0.2)
    TBL_T = Inches(0.5)
    TBL_W = Inches(12.9)
    TBL_H = Inches(6.6)
    HDR_H = Inches(0.4)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    def set_bg(slide, c):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = c

    def set_cell_bg(cell, rgb):
        tcPr = cell._tc.get_or_add_tcPr()
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'),
                                  {'val': f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'})
        sf.append(sc)
        tcPr.append(sf)

    def ct(cell, text, size=10, bold=False, color=DARK,
           align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = valign
        except:
            pass
        for a in ['margin_left', 'margin_right', 'margin_top', 'margin_bottom']:
            setattr(cell, a, Pt(4))
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = color
        r.font.bold = bold
        try:
            p.line_spacing = Pt(size * 1.25)
        except:
            pass

    def cm(cell, lines, size=9, color=DARK):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except:
            pass
        for a in ['margin_left', 'margin_right', 'margin_top', 'margin_bottom']:
            setattr(cell, a, Pt(4))
        for i, (label, text) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(1)
            p.space_after = Pt(1)
            r = p.add_run()
            r.text = f"{label}  {text}"
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.color.rgb = color
            r.font.bold = True
            try:
                p.line_spacing = Pt(size * 1.3)
            except:
                pass

    def add_img(slide, path, box_l, box_t, box_w, box_h):
        if not os.path.exists(path):
            return
        try:
            with Image.open(path) as img:
                iw, ih = img.size
        except:
            return
        m = Inches(0.04)
        bw = box_w - m * 2
        bh = box_h - m * 2
        scale = min(bw / 12700 / iw, bh / 12700 / ih)
        dw = int(iw * scale * 12700)
        dh = int(ih * scale * 12700)
        ox = int((bw - dw) / 2)
        oy = int((bh - dh) / 2)
        try:
            slide.shapes.add_picture(path, box_l + m + ox, box_t + m + oy, dw, dh)
        except:
            pass

    pages = [shots_data[i:i + 6] for i in range(0, len(shots_data), 6)]

    for pi, pd in enumerate(pages):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s, PAGE_BG)
        nr = len(pd) + 1
        ts = s.shapes.add_table(nr, ncols, TBL_L, TBL_T, TBL_W, TBL_H)
        tb = ts.table
        for ci, cw in enumerate(col_ws):
            tb.columns[ci].width = Inches(cw)

        data_h = Emu(int((TBL_H - HDR_H) / len(pd)))
        tb.rows[0].height = HDR_H
        for ri in range(1, nr):
            tb.rows[ri].height = data_h

        for ci, h in enumerate(col_hdrs):
            cell = tb.cell(0, ci)
            ct(cell, h, 10, True, WHITE)
            set_cell_bg(cell, (0x1A, 0x1A, 0x1A))

        row_ys = [TBL_T + HDR_H + i * data_h for i in range(len(pd))]

        for ri, sh in enumerate(pd):
            r_idx = ri + 1
            ry = row_ys[ri]
            bg = (0xF2, 0xF2, 0xF2) if ri % 2 == 0 else (0xFF, 0xFF, 0xFF)
            for ci in range(ncols):
                set_cell_bg(tb.cell(r_idx, ci), bg)

            ct(tb.cell(r_idx, 0), str(sh["num"]), 13, True, BLACK)
            ct(tb.cell(r_idx, 2), sh["type"], 10, True)
            ct(tb.cell(r_idx, 3), sh["time"], 10)
            cm(tb.cell(r_idx, 4), [
               ("场景", sh["scene"]),
               ("人物", sh["person"]),
               ("动作", sh["action"]),
            ], 9)
            ct(tb.cell(r_idx, 5), sh["script"], 9)
            ct(tb.cell(r_idx, 6), sh["camera"], 10, True)
            ct(tb.cell(r_idx, 7), sh["prop"], 10)

            ip = f"shots/shot_{sh['num']:02d}.jpg"
            ix = TBL_L + Inches(sum(col_ws[:1]))
            iw = Inches(col_ws[1])
            add_img(s, ip, ix, ry, iw, data_h)

        # 页眉
        tf = s.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(10), Inches(0.4)).text_frame
        r = tf.paragraphs[0].add_run()
        r.text = f"{title}  分镜头脚本  {pi + 1}/{len(pages)}"
        r.font.size = Pt(12)
        r.font.name = FONT
        r.font.color.rgb = BLACK
        r.font.bold = True

        tf2 = s.shapes.add_textbox(Inches(10.5), Inches(0.08), Inches(2.5), Inches(0.4)).text_frame
        tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
        r2 = tf2.paragraphs[0].add_run()
        r2.text = f"{len(shots_data)}镜"
        r2.font.size = Pt(9)
        r2.font.name = FONT
        r2.font.color.rgb = MID

        tf3 = s.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(8), Inches(0.25)).text_frame
        r3 = tf3.paragraphs[0].add_run()
        r3.text = "出品：大洋（乔治小弟）· 杭州短片导演"
        r3.font.size = Pt(7)
        r3.font.name = FONT
        r3.font.color.rgb = MID

    prs.save(output_path)
    print(f"✅ PPT 已生成: {output_path}")
    return len(prs.slides)


def cleanup():
    """清理临时文件"""
    for f in ["temp_video.mp4"]:
        if os.path.exists(f):
            os.remove(f)
    if os.path.exists("shots"):
        for f in os.listdir("shots"):
            os.remove(os.path.join("shots", f))
        os.rmdir("shots")


def main():
    if len(sys.argv) < 2:
        print("用法: python3 storyboard_generator.py <视频链接>")
        print("示例: python3 storyboard_generator.py https://www.bilibili.com/video/BVxxx")
        sys.exit(1)

    url = sys.argv[1]

    try:
        # Step 1: 下载
        meta = download_video(url)
        title = meta.get("title", "未命名视频").replace("/", "·")

        # Step 2: 场景检测
        scene_times, duration = detect_scenes()

        # Step 3: 构建镜头分段
        segments = build_shot_segments(scene_times, duration)
        print(f"📽️ 共 {len(segments)} 个镜头")

        # Step 4: 提取关键帧
        extract_frames(segments)

        # Step 5: 生成分镜数据
        title, shots = generate_storyboard_data(meta, segments)
        print(f"📝 分镜数据已生成，请编辑 {title}.py 中的 scene/person/action 字段")

        # Step 6: 生成 PPT
        safe_title = title.replace("/", "·").replace(" ", "_")
        output_path = os.path.join(OUTPUT_DIR, f"{safe_title}-分镜头脚本.pptx")
        make_pptx(shots, title, output_path)

        # Step 7: 清理
        cleanup()

        print(f"\n🎉 完成！文件已保存到: {output_path}")
        print("📌 提示：打开PPT后，请补充「景别」「运镜」「场景/人物/动作」字段")

    except Exception as e:
        print(f"❌ 错误: {e}")
        cleanup()
        sys.exit(1)


if __name__ == "__main__":
    main()
